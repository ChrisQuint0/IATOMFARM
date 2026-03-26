import io
import os
import tempfile
import fitz # PyMuPDF
import docx
from pptx import Presentation
from typing import List
from pydantic import BaseModel
from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import StreamingResponse, Response
from fastapi.middleware.cors import CORSMiddleware
from ai_utils import generate_flashcards
from export_utils import (
    create_csv_export, 
    create_docx_export, 
    create_pdf_export, 
    create_apkg_export
)

app = FastAPI(title="IATOMFARM API")

# Models for Export
class Flashcard(BaseModel):
    question: str
    answer: str

class ExportRequest(BaseModel):
    cards: List[Flashcard]
    deck_name: str = "IATOMFARM Deck"

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/health")
async def health_check():
    return {"status": "ok", "message": "IATOMFARM Backend is running"}

@app.post("/upload")
async def upload_document(
    file: UploadFile = File(...),
    card_count: int = Form(20),
    include_enumeration: bool = Form(False)
):
    filename = file.filename
    content_type = file.content_type
    
    try:
        file_bytes = await file.read()
        extracted_text = ""

        if filename.endswith(".pdf") or content_type == "application/pdf":
            # PDF Parsing with PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page in doc:
                extracted_text += page.get_text()
            doc.close()

        elif filename.endswith(".docx") or content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # DOCX Parsing with python-docx
            doc = docx.Document(io.BytesIO(file_bytes))
            extracted_text = "\n".join([para.text for para in doc.paragraphs])

        elif filename.endswith(".pptx") or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # PPTX Parsing with python-pptx
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"

        else:
            raise HTTPException(status_code=400, detail="Unsupported file format. Please upload a PDF, DOCX, or PPTX.")

        if not extracted_text.strip():
            raise HTTPException(status_code=422, detail="No text could be extracted from the document.")

        # Generate Flashcards using Gemini
        flashcards = generate_flashcards(
            text=extracted_text, 
            card_count=card_count, 
            include_enumeration=include_enumeration
        )

        return {
            "filename": filename,
            "total_cards": len(flashcards),
            "cards": flashcards
        }

    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")

@app.post("/export/{format}")
async def export_cards(format: str, request: ExportRequest):
    cards = [card.dict() for card in request.cards]
    
    if format == "csv":
        content = create_csv_export(cards)
        return Response(
            content=content,
            media_type="text/csv",
            headers={"Content-Disposition": "attachment; filename=iatomfarm_flashcards.csv"}
        )
    
    elif format == "docx":
        file_stream = create_docx_export(cards)
        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=iatomfarm_reviewer.docx"}
        )
    
    elif format == "pdf":
        pdf_bytes = create_pdf_export(cards)
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=iatomfarm_reviewer.pdf"}
        )
    
    elif format == "apkg":
        # APKG needs a temporary file because genanki writes to disk
        temp_dir = tempfile.gettempdir()
        temp_path = os.path.join(temp_dir, "iatomfarm_deck.apkg")
        
        import random
        import genanki
        
        model_id = random.randrange(1 << 30, 1 << 31)
        deck_id = random.randrange(1 << 30, 1 << 31)
        
        my_model = genanki.Model(
            model_id,
            'IATOMFARM Simple Model',
            fields=[{'name': 'Question'}, {'name': 'Answer'}],
            templates=[{
                'name': 'Card 1',
                'qfmt': '<div style="font-family: Arial; font-size: 20px; text-align: center; padding: 20px;">{{Question}}</div>',
                'afmt': '{{FrontSide}}<hr id="answer"><div style="font-family: Arial; font-size: 20px; text-align: center; color: #4F46E5; font-weight: bold; padding: 20px;">{{Answer}}</div>',
            }]
        )
        
        my_deck = genanki.Deck(deck_id, request.deck_name)
        for card in cards:
            my_note = genanki.Note(model=my_model, fields=[card['question'], card['answer']])
            my_deck.add_note(my_note)
            
        genanki.Package(my_deck).write_to_file(temp_path)
        
        with open(temp_path, "rb") as f:
            apkg_bytes = f.read()
            
        os.remove(temp_path)
        
        return Response(
            content=apkg_bytes,
            media_type="application/octet-stream",
            headers={"Content-Disposition": f"attachment; filename={request.deck_name.replace(' ', '_')}.apkg"}
        )

    else:
        raise HTTPException(status_code=400, detail="Unsupported export format.")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
